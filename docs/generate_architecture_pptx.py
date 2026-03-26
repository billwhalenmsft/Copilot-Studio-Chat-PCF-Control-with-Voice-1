"""
Generate 'The Microsoft Advantage' architecture diagram as an editable PowerPoint.
Recreated from the Copilot Studio Chat PCF Control - Voice & Vision diagram.
v2 — Major visual overhaul: proper vertical text via WordArt/stacked chars,
     distinct icon letters, better spacing, cleaner layout.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import os

# ── Colour palette ──────────────────────────────────────────────
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x00, 0x00, 0x00)
DARK_GREY    = RGBColor(0x33, 0x33, 0x33)
MED_GREY     = RGBColor(0x60, 0x60, 0x60)
LIGHT_GREY   = RGBColor(0xF2, 0xF2, 0xF2)
BORDER_GREY  = RGBColor(0xD0, 0xD0, 0xD0)
MS_BLUE      = RGBColor(0x00, 0x78, 0xD4)
AZURE_BLUE   = RGBColor(0x00, 0x89, 0xD6)
NAVY         = RGBColor(0x00, 0x33, 0x66)
D365_BLUE    = RGBColor(0x00, 0x2D, 0x62)
PP_PURPLE    = RGBColor(0x74, 0x2B, 0x8F)
COPILOT_PURPLE = RGBColor(0x7B, 0x2D, 0x8E)
TEAMS_PURPLE = RGBColor(0x50, 0x59, 0xC9)
GREEN        = RGBColor(0x10, 0x7C, 0x41)
SOLUTION_GRN = RGBColor(0x19, 0x8E, 0x3E)
TEAL         = RGBColor(0x00, 0x80, 0x80)
ORANGE       = RGBColor(0xE3, 0x6C, 0x09)
GOLD         = RGBColor(0xF2, 0xC8, 0x11)
DV_GREEN     = RGBColor(0x00, 0x6B, 0x3F)

# Brand-accurate icon colours
IC = {
    "pa_mobile": PP_PURPLE,
    "web":       MS_BLUE,
    "mda":       RGBColor(0x95, 0x4F, 0x72),
    "teams":     TEAMS_PURPLE,
    "outlook":   MS_BLUE,
    "portals":   MS_BLUE,
    "tools":     RGBColor(0x10, 0x7C, 0x10),
    "bi_3p":     GOLD,
    "custom":    TEAL,
    "d365":      D365_BLUE,
    "pp_apps":   PP_PURPLE,
    "pp_auto":   RGBColor(0x00, 0x66, 0xFF),
    "pp_va":     RGBColor(0x0B, 0x55, 0x6A),
    "pp_pages":  RGBColor(0x8C, 0x4F, 0xFF),
    "pp_bi":     GOLD,
    "copilot":   COPILOT_PURPLE,
    "graph":     MS_BLUE,
    "dv":        DV_GREEN,
}

# ── Helpers ─────────────────────────────────────────────────────

def _tb(slide, left, top, width, height):
    """Raw textbox shape."""
    return slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))


def add_text(slide, left, top, width, height, text, size=10,
             bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             font="Segoe UI", wrap=True, anchor=MSO_ANCHOR.TOP):
    """General-purpose textbox."""
    tb = _tb(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = align
    except Exception:
        pass
    # Set vertical anchor
    bodyPr = tf._txBody.bodyPr
    anchor_map = {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}
    bodyPr.set("anchor", anchor_map.get(anchor, "t"))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = Pt(size + 2)
    return tb


def add_rich_text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT, wrap=True):
    """Textbox with multiple styled runs.  runs = [(text, size, bold, color), …]"""
    tb = _tb(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    for i, (txt, sz, bld, clr) in enumerate(runs):
        r = p.add_run() if i > 0 else p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bld
        r.font.color.rgb = clr
        r.font.name = "Segoe UI"
    # Remove the auto-created empty first run that paragraph starts with
    if len(p.runs) > 0 and p.runs[0].text == "":
        p._p.remove(p.runs[0]._r)
    return tb


def add_rect(slide, left, top, width, height, fill=None,
             border=None, bw=Pt(1), rounded=False):
    """Rectangle or rounded rectangle."""
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp_type,
        Inches(left), Inches(top), Inches(width), Inches(height))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = bw
    else:
        s.line.fill.background()
    # Reduce default padding
    tf = s.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    return s


def add_pill(slide, left, top, width, height, text, fill_color,
             text_color=WHITE, size=7.5, bold=True):
    """Rounded badge / pill."""
    s = add_rect(slide, left, top, width, height,
                 fill=fill_color, rounded=True)
    tf = s.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    bodyPr = tf._txBody.bodyPr
    bodyPr.set("anchor", "ctr")
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    return s


def add_circle(slide, left, top, diameter, fill_color, letter,
               letter_size=None, letter_color=WHITE):
    """Coloured circle with a single centred letter."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(diameter), Inches(diameter))
    c.fill.solid()
    c.fill.fore_color.rgb = fill_color
    c.line.fill.background()
    tf = c.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    bodyPr = tf._txBody.bodyPr
    bodyPr.set("anchor", "ctr")
    lsz = letter_size or max(int(diameter * 22), 8)
    p = tf.paragraphs[0]
    p.text = letter
    p.font.size = Pt(lsz)
    p.font.bold = True
    p.font.color.rgb = letter_color
    p.font.name = "Segoe UI Semibold"
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    return c


def icon_with_label(slide, cx, cy, d, color, letter, label,
                    label_lines=1, lbl_size=6):
    """Circle icon + centred label below. Returns (circle, label_tb)."""
    c = add_circle(slide, cx, cy, d, color, letter)
    lw = max(d + 0.3, 0.55)
    lh = 0.14 * label_lines + 0.06
    tb = add_text(slide, cx - (lw - d) / 2, cy + d + 0.02,
                  lw, lh, label, size=lbl_size, color=DARK_GREY,
                  align=PP_ALIGN.CENTER, wrap=True)
    return c, tb


def vertical_bar(slide, left, top, width, height, fill_color, label_chars):
    """Vertical side-bar with stacked characters (one per line) for reliable
    vertical text rendering — avoids rotated textbox rendering issues."""
    add_rect(slide, left, top, width, height, fill=fill_color)
    # Stack characters vertically using line breaks
    stacked = "\n".join(label_chars)
    tb = _tb(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    bodyPr = tf._txBody.bodyPr
    bodyPr.set("anchor", "ctr")
    p = tf.paragraphs[0]
    p.text = stacked
    p.font.size = Pt(5)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = Pt(7)
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    return tb


def bullet_row(slide, left, top, width, bold_part, rest_part, size=8.5):
    """Bullet point: • Bold prefix + normal rest.  Returns textbox."""
    tb = _tb(slide, left, top, width, 0.5)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = "• "
    r1.font.size = Pt(size)
    r1.font.color.rgb = DARK_GREY
    r1.font.name = "Segoe UI"
    r2 = p.add_run()
    r2.text = bold_part
    r2.font.size = Pt(size)
    r2.font.bold = True
    r2.font.color.rgb = DARK_GREY
    r2.font.name = "Segoe UI"
    r3 = p.add_run()
    r3.text = rest_part
    r3.font.size = Pt(size)
    r3.font.bold = False
    r3.font.color.rgb = DARK_GREY
    r3.font.name = "Segoe UI"
    return tb


# ── Slide builder ───────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # White background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # ── Coordinate constants ──
    DL = 0.5            # diagram left
    DT = 1.05           # diagram top
    DW = 8.2            # diagram width
    DH = 5.95           # diagram height
    BAR_W = 0.32        # side-bar width
    IX = DL + BAR_W + 0.12   # inner content left
    IW = DW - 2 * (BAR_W + 0.12)  # inner content width
    IR = IX + IW              # inner content right

    # ================================================================
    #  TITLE
    # ================================================================
    add_text(slide, 0.5, 0.15, 7, 0.4,
             "Copilot Studio Chat PCF Control",
             size=20, bold=True, color=DARK_GREY)
    add_text(slide, 0.5, 0.52, 7, 0.35,
             "Voice & Vision — The Microsoft Advantage",
             size=15, bold=True, color=MS_BLUE)
    add_text(slide, 11.5, 0.2, 1.6, 0.35, "Microsoft",
             size=14, bold=True, color=DARK_GREY, align=PP_ALIGN.RIGHT,
             font="Segoe UI Semibold")

    # ================================================================
    #  OUTER FRAME + SIDE BARS
    # ================================================================
    add_rect(slide, DL, DT, DW, DH, border=BORDER_GREY, bw=Pt(1.5))

    # Left bar: SECURITY, COMPLIANCE & IDENTITY
    vertical_bar(slide, DL, DT, BAR_W, DH, NAVY,
                 list("SECURITY, COMPLIANCE & IDENTITY"))
    # Right bar: MANAGEMENT & GOVERNANCE
    vertical_bar(slide, DL + DW - BAR_W, DT, BAR_W, DH, NAVY,
                 list("MANAGEMENT & GOVERNANCE"))

    # ================================================================
    #  ROW 1 — User Interface & Mobile  |  3rd-party Applications
    # ================================================================
    r1y = DT + 0.12
    # UI box
    ui_w = IW * 0.46
    add_rect(slide, IX, r1y, ui_w, 1.35,
             fill=LIGHT_GREY, border=BORDER_GREY)
    add_text(slide, IX + 0.12, r1y + 0.06, ui_w, 0.22,
             "User Interface & Mobile", size=9, bold=True, color=DARK_GREY)

    ui_icons = [
        ("P", "Power Apps\nMobile", IC["pa_mobile"]),
        ("W", "Web\nBrowser",       IC["web"]),
        ("M", "Model\nDriven App",  IC["mda"]),
        ("T", "Teams\nChat",        IC["teams"]),
    ]
    ui_x0 = IX + 0.18
    for i, (ltr, lbl, clr) in enumerate(ui_icons):
        icon_with_label(slide, ui_x0 + i * 0.78, r1y + 0.32, 0.42,
                        clr, ltr, lbl, label_lines=2, lbl_size=5.5)

    # Outlook below row
    icon_with_label(slide, ui_x0, r1y + 0.95, 0.35,
                    IC["outlook"], "O", "Outlook", lbl_size=5.5)

    # 3rd-party box
    tp_x = IX + ui_w + 0.15
    tp_w = IW - ui_w - 0.15
    add_rect(slide, tp_x, r1y, tp_w, 0.95,
             fill=LIGHT_GREY, border=BORDER_GREY)
    add_text(slide, tp_x + 0.12, r1y + 0.06, tp_w, 0.22,
             "3rd-party Applications", size=9, bold=True, color=DARK_GREY)

    tp_icons = [
        ("W", "Web\nPortals",  IC["portals"]),
        ("T", "Tools",         IC["tools"]),
        ("B", "BI",            IC["bi_3p"]),
        ("C", "Custom\nIntegrations", IC["custom"]),
    ]
    tp_x0 = tp_x + 0.2
    for i, (ltr, lbl, clr) in enumerate(tp_icons):
        icon_with_label(slide, tp_x0 + i * 0.82, r1y + 0.3, 0.4,
                        clr, ltr, lbl, label_lines=2, lbl_size=5.5)

    # "3rd-party Data" label right-aligned inside the 3rd-party box
    add_text(slide, tp_x + tp_w - 1.2, r1y + 0.72, 1.1, 0.2,
             "3rd-party Data", size=7.5, color=MED_GREY, align=PP_ALIGN.RIGHT)

    # ================================================================
    #  ROW 2 — Microsoft Graph  …  AI Copilot
    # ================================================================
    r2y = r1y + 1.5
    add_text(slide, IX, r2y + 0.08, 1.8, 0.22,
             "Microsoft Graph", size=9, bold=True, color=DARK_GREY)

    # Copilot circle (centered)
    cop_d = 0.62
    cop_cx = IX + IW / 2 - cop_d / 2
    cop = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(cop_cx), Inches(r2y - 0.05),
        Inches(cop_d), Inches(cop_d))
    cop.fill.solid()
    cop.fill.fore_color.rgb = WHITE
    cop.line.color.rgb = COPILOT_PURPLE
    cop.line.width = Pt(3)
    # small copilot icon text inside
    tf = cop.text_frame
    tf.word_wrap = False
    bodyPr = tf._txBody.bodyPr
    bodyPr.set("anchor", "ctr")
    p = tf.paragraphs[0]
    p.text = "✦"
    p.font.size = Pt(16)
    p.font.color.rgb = COPILOT_PURPLE
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    add_text(slide, cop_cx - 0.15, r2y + cop_d - 0.02, cop_d + 0.3, 0.2,
             "AI Copilot", size=7, bold=True, color=COPILOT_PURPLE,
             align=PP_ALIGN.CENTER)

    # ================================================================
    #  ROW 3 — Unified CRM  |  Low Code
    # ================================================================
    r3y = r2y + 0.65
    col_gap = 0.2
    crm_w = IW * 0.46
    lc_w  = IW - crm_w - col_gap
    crm_h = 2.4
    lc_x  = IX + crm_w + col_gap

    # ── Unified CRM ──
    add_rect(slide, IX, r3y, crm_w, crm_h, border=BORDER_GREY)
    add_text(slide, IX + 0.12, r3y + 0.06, 1.5, 0.22,
             "Unified CRM", size=10, bold=True, color=DARK_GREY)

    # D365 header
    d365y = r3y + 0.32
    add_circle(slide, IX + 0.12, d365y, 0.32, D365_BLUE, "D", letter_size=12)
    add_text(slide, IX + 0.52, d365y + 0.03, 1.8, 0.25,
             "Dynamics 365", size=10, bold=True, color=D365_BLUE)

    # D365 app icons
    d365_apps = [
        ("C", "CSS",               D365_BLUE),
        ("S", "Sales",             RGBColor(0x00, 0x78, 0xD4)),
        ("M", "Marketing &\nCust. Insights", RGBColor(0x9B, 0x59, 0xB6)),
        ("C", "Customer\nService", RGBColor(0x00, 0x95, 0x5E)),
        ("F", "Field\nService",    RGBColor(0x2E, 0x86, 0xC1)),
    ]
    ax0 = IX + 0.12
    for i, (ltr, lbl, clr) in enumerate(d365_apps):
        icon_with_label(slide, ax0 + i * 0.62, d365y + 0.38, 0.33,
                        clr, ltr, lbl, label_lines=2, lbl_size=4.5)

    # Copilot Studio
    csy = d365y + 1.1
    add_circle(slide, IX + 0.12, csy, 0.32, COPILOT_PURPLE, "C", letter_size=12)
    add_text(slide, IX + 0.52, csy + 0.02, 2.5, 0.25,
             "Microsoft Copilot Studio", size=9, bold=True, color=COPILOT_PURPLE)
    add_text(slide, IX + 0.12, csy + 0.35, crm_w - 0.24, 0.2,
             "Agent • Topics • Knowledge • Generative AI • Direct Line API v3",
             size=6.5, color=MED_GREY)

    # ── Low Code ──
    add_rect(slide, lc_x, r3y, lc_w, crm_h, border=BORDER_GREY)
    add_text(slide, lc_x + 0.12, r3y + 0.06, 1.5, 0.22,
             "Low Code", size=10, bold=True, color=DARK_GREY)

    # Power Platform sub-box
    pp_y = r3y + 0.32
    pp_h = 0.92
    add_rect(slide, lc_x + 0.1, pp_y, lc_w - 0.2, pp_h,
             fill=LIGHT_GREY, border=BORDER_GREY)
    add_text(slide, lc_x + 0.2, pp_y + 0.04, 2.0, 0.2,
             "Power Platform", size=9, bold=True, color=PP_PURPLE)

    pp_icons = [
        ("PA", "Power\nApps",       IC["pp_apps"]),
        ("PF", "Power\nAutomate",   IC["pp_auto"]),
        ("PV", "Power Virtual\nAgents", IC["pp_va"]),
        ("PP", "Power\nPages",      IC["pp_pages"]),
    ]
    ppx0 = lc_x + 0.22
    for i, (ltr, lbl, clr) in enumerate(pp_icons):
        icon_with_label(slide, ppx0 + i * 0.78, pp_y + 0.28, 0.35,
                        clr, ltr, lbl, label_lines=2, lbl_size=4.5)

    # Power BI mini icon
    icon_with_label(slide, ppx0, pp_y + 0.68, 0.22,
                    IC["pp_bi"], "BI", "Power BI", lbl_size=4.5)

    # ── OUR SOLUTION — PCF Chat Controls ──
    pcf_y = pp_y + pp_h + 0.12
    pcf_h = 1.0
    add_rect(slide, lc_x + 0.1, pcf_y, lc_w - 0.2, pcf_h,
             border=SOLUTION_GRN, bw=Pt(2.5), rounded=True)

    add_pill(slide, lc_x + 0.18, pcf_y + 0.06, 1.1, 0.2,
             "◆ OUR SOLUTION", SOLUTION_GRN, WHITE, size=6.5)

    add_text(slide, lc_x + 0.18, pcf_y + 0.28, lc_w - 0.4, 0.2,
             "PCF Chat Controls", size=10, bold=True, color=DARK_GREY)
    add_text(slide, lc_x + 0.18, pcf_y + 0.48, lc_w - 0.4, 0.18,
             "CopilotChatBeta v2.0.4  |  CopilotChatGA  |  CopilotVisionChat",
             size=7, color=MED_GREY)

    # Feature badges — row 1
    badges1 = [
        ("Audio",          MS_BLUE),
        ("Vision",         TEAL),
        ("Attachments",    ORANGE),
        ("34 Languages",   GREEN),
    ]
    bx = lc_x + 0.18
    for lbl, clr in badges1:
        bw = len(lbl) * 0.06 + 0.22
        add_pill(slide, bx, pcf_y + 0.67, bw, 0.19, lbl, clr, WHITE, size=6.5)
        bx += bw + 0.06

    # Feature badges — row 2
    badges2 = [
        ("Driving Mode",   MS_BLUE),
        ("Adaptive Cards",  ORANGE),
    ]
    bx2 = lc_x + 0.18
    for lbl, clr in badges2:
        bw = len(lbl) * 0.06 + 0.22
        add_pill(slide, bx2, pcf_y + 0.87 if pcf_h > 0.95 else pcf_y + 0.78,
                 bw, 0.19, lbl, clr, WHITE, size=6.5)
        bx2 += bw + 0.06

    # ================================================================
    #  ROW 4 — Microsoft Graph & WorkIQ
    # ================================================================
    r4y = r3y + crm_h + 0.12
    add_rect(slide, IX, r4y, IW, 0.38,
             fill=LIGHT_GREY, border=BORDER_GREY)
    add_circle(slide, IX + 0.08, r4y + 0.04, 0.3,
               IC["graph"], "G", letter_size=11)
    add_text(slide, IX + 0.46, r4y + 0.03, 2.5, 0.18,
             "Microsoft Graph & WorkIQ", size=9, bold=True, color=DARK_GREY)
    add_text(slide, IX + 0.46, r4y + 0.2, IW - 0.6, 0.16,
             "Unified API for people, files, mail, calendar, tasks • Intelligent work insights & analytics",
             size=6.5, color=MED_GREY)

    # ================================================================
    #  ROW 5 — Dataverse
    # ================================================================
    r5y = r4y + 0.46
    add_rect(slide, IX, r5y, IW, 0.32,
             fill=LIGHT_GREY, border=BORDER_GREY)
    add_circle(slide, IX + 0.08, r5y + 0.03, 0.26,
               IC["dv"], "D", letter_size=10)
    add_text(slide, IX + 0.42, r5y + 0.04, 1.2, 0.18,
             "Dataverse", size=9, bold=True, color=DARK_GREY)
    add_text(slide, IX + 1.25, r5y + 0.04, IW - 1.4, 0.18,
             "Configuration • Environment Variables • Entity Data • Solutions",
             size=7, color=MED_GREY)

    # ================================================================
    #  ROW 6 — Azure AI
    # ================================================================
    r6y = r5y + 0.4
    az_h = 0.6
    add_rect(slide, IX, r6y, IW, az_h, fill=NAVY)
    add_text(slide, IX + 0.15, r6y + 0.04, 1.3, 0.25,
             "Azure AI", size=12, bold=True, color=WHITE)
    add_text(slide, IX + 0.15, r6y + 0.28, 2.8, 0.28,
             "Global cloud infrastructure & services\nIdentity, security, management, and compliance",
             size=6, color=RGBColor(0xBB, 0xD5, 0xEE))

    # Service pills
    az_svc = [
        ("Speech",         AZURE_BLUE),
        ("OpenAI",         RGBColor(0x10, 0xA3, 0x7F)),
        ("Vision Proxy",   TEAL),
        ("Token Exchange",  ORANGE),
        ("Entra ID",       RGBColor(0x00, 0xA4, 0xEF)),
        ("Functions",      RGBColor(0xE8, 0xA8, 0x00)),
    ]
    abx = IX + 2.6
    for lbl, clr in az_svc:
        bw = len(lbl) * 0.065 + 0.25
        add_pill(slide, abx, r6y + 0.18, bw, 0.24,
                 lbl, clr, WHITE, size=7.5, bold=True)
        abx += bw + 0.1

    # ================================================================
    #  RIGHT COLUMN — Value Propositions
    # ================================================================
    VP_X = DL + DW + 0.3
    VP_W = 13.333 - VP_X - 0.3

    vps = [
        ("A complete, unified suite ",
         "of leading front and back-office business applications with embedded AI chat"),
        ("The #1 platform ",
         "for creating personal and enterprise-wide automations and low code apps"),
        ("Single common data model ",
         "with AI built-in every product via Copilot Studio agents"),
        ("Seamless integration ",
         "across Microsoft — D365, Power Platform, Graph, and Azure AI services"),
        ("Built on Azure AI ",
         "with enterprise security, compliance, privacy and Responsible AI"),
        ("Voice, Vision & 34 languages — ",
         "PCF controls add speech, camera, and multilingual chat to any Power App"),
        ("Microsoft Graph & WorkIQ ",
         "surface people, files, calendar, and intelligent work insights directly in the chat experience"),
        ("Outlook & Teams Chat integration ",
         "extends Copilot conversations into everyday productivity workflows"),
        ("End-to-end security — ",
         "Entra ID, RBAC, conditional access, compliance, and governance scaffold every layer of the stack"),
    ]

    bpy = DT + 0.05
    for bold_part, rest_part in vps:
        bullet_row(slide, VP_X, bpy, VP_W, bold_part, rest_part, size=8.5)
        bpy += 0.6

    return prs


# ── Generate ────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = build()
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "CopilotStudioChat-Architecture-Diagram-v2.pptx")
    prs.save(out)
    print(f"✅  Saved → {out}")
